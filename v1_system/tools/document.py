"""document.py - 多模态文档生成系统"""
from pathlib import Path
from typing import Optional, List, Dict, Any

class WordDocument:
    def __init__(self, path: str):
        try:
            from docx import Document
        except ImportError:
            raise ImportError("pip install python-docx")
        self.path = Path(path)
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.doc = Document(str(self.path)) if self.path.exists() else Document()
    
    def add_heading(self, text: str, level: int = 1) -> None:
        self.doc.add_heading(text, level=level)
    
    def add_paragraph(self, text: str) -> None:
        self.doc.add_paragraph(text)
    
    def add_bullet_list(self, items: List[str]) -> None:
        for item in items:
            self.doc.add_paragraph(item, style='List Bullet')
    
    def save(self) -> bool:
        try:
            self.doc.save(str(self.path))
            return True
        except Exception as e:
            print(f"Save failed: {e}")
            return False

class PPTPresentation:
    def __init__(self, path: str):
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("pip install python-pptx")
        self.path = Path(path)
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.prs = Presentation(str(self.path)) if self.path.exists() else Presentation()
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
    
    def add_content_slide(self, title: str, content: List[str]) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = title
        for item in content:
            p = slide.placeholders[1].text_frame.add_paragraph()
            p.text = item
    
    def save(self) -> bool:
        try:
            self.prs.save(str(self.path))
            return True
        except Exception as e:
            print(f"Save failed: {e}")
            return False
