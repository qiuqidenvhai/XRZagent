"""
commander.py — 仙人掌 Agent 主控制器
- 总工程师：调度工具 + 子代理 + 记忆
- 支持 continue/remember/recall 指令
- 自动摘要提醒
"""
import asyncio
import time
import logging
import uuid
import re
import os
import sys
from dataclasses import dataclass
from enum import Enum
from typing import Optional, Callable
from types import SimpleNamespace

from .protocol import Protocol, ExecutionResult
from .session import DeepSeekSession, SessionConfig, MessageRole
from .subagent import SubAgentManager, BrowserSubAgent, SubAgentResult
from .memory_manager import MemoryManager

logger = logging.getLogger("commander")


class EventType(Enum):
    THINKING = "thinking"
    TOOL_START = "tool_start"
    TOOL_END = "tool_end"
    TOOL_ERROR = "tool_error"
    COMMAND_DETECTED = "command_detected"
    COMMAND_EXECUTING = "command_executing"
    COMMAND_SUCCESS = "command_success"
    COMMAND_ERROR = "command_error"
    AI_FINAL_REPLY = "ai_final_reply"
    MEMORY_REMINDER = "memory_reminder"
    CONTINUE_READY = "continue_ready"
    CORRECTION_SENT = "correction_sent"
    ERROR = "error"


@dataclass
class AgentEvent:
    event_type: EventType
    data: any = None


# ============================================================
# 工具注册表
# ============================================================
class ToolRegistry:
    """工具注册表，支持子代理工具"""

    def __init__(self, commander):
        self._commander = commander
        self._tools: dict = {}
        self._browser_subagent: Optional[BrowserSubAgent] = None

    def set_browser_subagent(self, bm):
        """设置浏览器子代理（兼容旧接口，不再使用）"""
        self._browser_subagent = BrowserSubAgent(bm, depth=1)

    def register(self, name: str, description: str, fn: Callable):
        self._tools[name] = {"description": description, "fn": fn}

    def list_tools(self) -> list:
        return [
            {"name": name, "description": info["description"]}
            for name, info in self._tools.items()
        ]

    async def execute(self, tool_name: str, params: dict) -> "ExecutionResult":
        """异步执行工具（由 Commander 调用）"""
        import uuid
        if tool_name not in self._tools:
            return ExecutionResult(id=str(uuid.uuid4()), status="error",
                                  error=f"未知工具: {tool_name}", tool=tool_name)
        fn = self._tools[tool_name]["fn"]
        try:
            result = await fn(**params) if asyncio.iscoroutinefunction(fn) else fn(**params)
            return ExecutionResult(id=str(uuid.uuid4()), status="success", output=str(result), tool=tool_name)
        except Exception as e:
            return ExecutionResult(id=str(uuid.uuid4()), status="error", error=str(e), tool=tool_name)


# ============================================================
# 主控制器 Commander
# ============================================================

class Commander:
    """
    Agent 主控制器（总工程师模式）
    - 对话管理：维护 session + 事件回调
    - 工具执行：调用工具注册表
    - 子代理管理：维护 SubAgentManager（任务状态 + 凭据共享）
    - 记忆管理：自动摘要 + 长期记忆（remember/recall）
    - 错误纠正：commander fix
    """

    def __init__(
        self,
        browser_manager,
        session: Optional[DeepSeekSession] = None,
        work_dir: str = "",
        on_event: Optional[Callable[[AgentEvent], None]] = None,
    ):
        self._bm = browser_manager
        self._session = session
        self._work_dir = work_dir or os.getcwd()
        self._on_event = on_event
        self._tools = ToolRegistry(self)
        self._subagent_manager: Optional[SubAgentManager] = None
        self._memory: Optional[MemoryManager] = None
        self._history_turns: int = 0
        self._accumulated_prompt: str = ""
        self._correction_pending: Optional[str] = None
        self._running = False
        self._protocol = Protocol()

        # 先注册工具，再构建系统提示词（提示词依赖工具列表）
        self._register_tools()
        self._system_prompt = self._build_system_prompt()

    def _register_tools(self):
        """注册所有工具（包括子代理工具）"""
        import os
        from pathlib import Path

        # ─── 文件操作工具 ───
        async def file_write(**params):
            p = Path(params.get("path", ""))
            content = params.get("content", "")
            full_path = Path(self._work_dir) / p if not p.is_absolute() else p
            full_path.parent.mkdir(parents=True, exist_ok=True)
            full_path.write_text(content, encoding="utf-8")
            return f"文件已写入: {full_path}"

        async def file_read(**params):
            p = Path(params.get("path", ""))
            full_path = Path(self._work_dir) / p if not p.is_absolute() else p
            if not full_path.exists():
                return f"错误: 文件不存在 {full_path}"
            return full_path.read_text(encoding="utf-8")

        async def file_list(**params):
            p = Path(params.get("path", "."))
            full_path = Path(self._work_dir) / p if not p.is_absolute() else p
            if not full_path.exists():
                return "目录不存在"
            items = list(full_path.iterdir())
            return "\n".join([f"{'[DIR]' if i.is_dir() else '[FILE]'} {i.name}" for i in items])

        async def dir_create(**params):
            p = Path(params.get("path", ""))
            full_path = Path(self._work_dir) / p if not p.is_absolute() else p
            full_path.mkdir(parents=True, exist_ok=True)
            return f"目录已创建: {full_path}"

        async def file_delete(**params):
            p = Path(params.get("path", ""))
            full_path = Path(self._work_dir) / p if not p.is_absolute() else p
            if full_path.exists():
                if full_path.is_file():
                    full_path.unlink()
                else:
                    import shutil
                    shutil.rmtree(full_path)
                return f"已删除: {full_path}"
            return f"文件不存在: {full_path}"

        self._tools.register("file_write", "写入文件", file_write)
        self._tools.register("file_read", "读取文件", file_read)
        self._tools.register("file_list", "列出目录", file_list)
        self._tools.register("dir_create", "创建目录", dir_create)
        self._tools.register("file_delete", "删除文件/目录", file_delete)

        # ─── Shell 执行工具 ───
        async def shell_exec(**params):
            import asyncio
            cmd = params.get("command", "")
            timeout = params.get("timeout", 60)
            try:
                proc = await asyncio.create_subprocess_shell(
                    cmd,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                )
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
                return f"[退出码 {proc.returncode}]\n{stdout.decode('utf-8', errors='ignore')}\n{stderr.decode('utf-8', errors='ignore')}"
            except asyncio.TimeoutError:
                proc.kill()
                return f"命令超时（>{timeout}s）"
            except Exception as e:
                return f"执行错误: {e}"

        self._tools.register("shell_exec", "执行Shell命令", shell_exec)

        # ─── 浏览器操作工具（母代理直接执行）───

        async def browser_click(**params):
            return await self._bm.browser_click(params.get("selector", ""))

        async def browser_fill(**params):
            return await self._bm.browser_fill(params.get("selector", ""), params.get("text", ""))

        async def browser_screenshot(**params):
            path = params.get("path", "screenshot.png")
            full_path = Path(self._work_dir) / path
            return await self._bm.browser_screenshot(str(full_path))

        async def browser_search(**params):
            """网页搜索工具 - 不使用浏览器，直接用 HTTP 请求抓取页面
            
            通过 Bing 搜索获取结果，抓取页面内容后返回给 AI。
            不会占用母代理的浏览器状态。
            """
            import subprocess
            import json
            from pathlib import Path
            
            query = params.get("query", "")
            max_pages = params.get("max_pages", 3)
            output_file = params.get("output_file", "")
            
            if not query:
                return "[错误] 缺少 query 参数"
            
            # 确定输出文件路径
            if not output_file:
                output_file = str(Path(self._work_dir) / "search_results.json")
            
            # 调用网页搜索脚本
            script_path = Path(__file__).parent.parent / "web_searcher.py"
            
            try:
                result = subprocess.run(
                    [sys.executable, str(script_path), query, str(max_pages), output_file],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                if result.returncode == 0:
                    # 读取结果文件
                    if Path(output_file).exists():
                        with open(output_file, "r", encoding="utf-8") as f:
                            data = json.load(f)
                        return f"[搜索完成] 抓取 {data.get('scraped_count', 0)} 个页面\n\n{data.get('findings', '')}"
                    else:
                        return result.stdout
                else:
                    return f"[搜索错误] {result.stderr}"
                    
            except subprocess.TimeoutExpired:
                return "[错误] 搜索超时（60秒）"
            except Exception as e:
                return f"[错误] {str(e)}"

        # ─── 深度思考工具（母代理执行，AI 调用）───
        async def _deep_think_tool(**params):
            """深度思考工具 - 由 AI 调用，实际控制浏览器按钮
            
            当 AI 认为当前问题需要深度思考时，调用此工具开启 DeepSeek 的深度思考模式。
            系统会在 AI 输出完成后自动关闭此模式。
            """
            enable_str = params.get("enable", "true")
            enable = enable_str.lower() in ("true", "1", "yes", "on")
            await self._session.set_deep_think(enable)
            return f"深度思考模式已{'开启' if enable else '关闭'}"

        self._tools.register("browser_click", "点击元素", browser_click)
        self._tools.register("browser_fill", "填写输入框", browser_fill)
        self._tools.register("browser_screenshot", "截图", browser_screenshot)
        self._tools.register("browser_search", "搜索", browser_search)
        self._tools.register("deep_think", "深度思考：AI认为需要深度思考时调用此指令开启DeepSeek深度思考模式，完成后自动关闭", _deep_think_tool)

        # ─── 研究代理工具（独立子代理进程）───
        async def _browser_research(**params):
            """启动研究子代理（独立进程）"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)
            
            query = params.get("query", "")
            max_pages = params.get("max_pages", 5)
            
            # 构建完整查询
            full_query = f"深度研究: {query} (最多访问 {max_pages} 个页面)"
            
            # 启动子代理进程（非阻塞）
            task_id = await sam.spawn_subagent(full_query, task_type="research")
            
            return f"[子代理已启动] task_id={task_id}\n查询: {query[:50]}...\n使用 check_task('{task_id}') 查看进度"

        async def _browser_visit(**params):
            """启动访问子代理（独立进程）"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)
            
            url = params.get("url", "")
            
            # 构建完整查询
            full_query = f"访问并分析网页: {url}"
            
            # 启动子代理进程（非阻塞）
            task_id = await sam.spawn_subagent(full_query, task_type="visit")
            
            return f"[子代理已启动] task_id={task_id}\nURL: {url}\n使用 check_task('{task_id}') 查看进度"
        
        async def _check_task(**params):
            """检查子代理任务状态"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)
            
            task_id = params.get("task_id", "")
            task = sam.check_task(task_id)
            
            if not task:
                return f"任务 {task_id} 不存在"
            
            status_info = f"任务: {task_id}\n状态: {task.status.value}\n类型: {task.task_type}\n查询: {task.query[:60]}"
            
            if task.result:
                status_info += f"\n结果: {'成功' if task.result.success else '失败'}"
                if task.result.output:
                    preview = task.result.output[:200] + "..." if len(task.result.output) > 200 else task.result.output
                    status_info += f"\n输出预览: {preview}"
                if task.result.files:
                    status_info += f"\n文件数: {len(task.result.files)}"
            
            return status_info
        
        async def _wait_task(**params):
            """等待子代理任务完成"""
            from .subagent_manager import get_subagent_manager
            sam = get_subagent_manager(self._work_dir)
            
            task_id = params.get("task_id", "")
            timeout = params.get("timeout", 300)
            
            task = await sam.wait_task(task_id, timeout=timeout)
            
            if not task:
                return f"等待任务 {task_id} 超时或任务不存在"
            
            if task.result and task.result.success:
                return f"任务 {task_id} 完成！\n输出: {task.result.output[:500]}"
            else:
                error = task.result.error if task.result else "未知错误"
                return f"任务 {task_id} 失败: {error}"

        self._tools.register("browser_research", "启动研究子代理（独立进程）", _browser_research)
        self._tools.register("browser_visit", "启动访问子代理（独立进程）", _browser_visit)
        self._tools.register("check_task", "检查子代理任务状态", _check_task)
        self._tools.register("wait_task", "等待子代理任务完成", _wait_task)

        # ─── Word 文档生成工具（母代理直接执行）───
        async def docx_tool_fn(**params):
            import os
            from pathlib import Path
            content_text = params.get("content", "")
            filename = params.get("filename", "report.docx")
            path = params.get("path", "")
            # 展开 %USERPROFILE% 等环境变量
            filename = os.path.expandvars(filename)
            if path:
                path = os.path.expandvars(path)
            # 如果指定了 path 参数，使用 path 作为目录
            if path:
                out_path = Path(path) / filename
            else:
                out_path = Path(self._work_dir) / filename
            out_path.parent.mkdir(parents=True, exist_ok=True)
            try:
                from docx import Document
                doc = Document()
                doc.add_heading(filename.replace(".docx", ""), 0)
                for para in content_text.split("\n"):
                    if para.strip():
                        doc.add_paragraph(para)
                doc.save(str(out_path))
                return f"Word 文档已生成: {out_path}"
            except ImportError:
                # 降级为 txt
                txt_path = out_path.with_suffix(".txt")
                txt_path.write_text(content_text, encoding="utf-8")
                return f"python-docx 未安装，生成文本文件: {txt_path}"

        self._tools.register("docx_create", "生成 Word 文档（python-docx）", docx_tool_fn)

        # ─── PPT 生成工具 ───
        async def pptx_tool_fn(**params):
            content_text = params.get("content", "")
            filename = params.get("filename", "slides.pptx")
            out_path = Path(self._work_dir) / filename
            out_path.parent.mkdir(parents=True, exist_ok=True)
            try:
                from pptx import Presentation
                prs = Presentation()
                lines = content_text.split("\n")
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = filename.replace(".pptx", "")
                body = slide.shapes.placeholders[1]
                tf = body.text_frame
                for line in lines[:10]:
                    p = tf.add_paragraph()
                    p.text = line
                prs.save(str(out_path))
                return f"PPT 已生成: {out_path}"
            except ImportError:
                return "python-pptx 未安装，无法生成 PPT"

        self._tools.register("pptx_create", "生成 PPT（python-pptx）", pptx_tool_fn)

        # ─── 记忆工具（母代理直接执行）───
        async def remember_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            content = params.get("content", "")
            tags = params.get("tags", [])
            mid = self._memory.save(content, tags=tags if isinstance(tags, list) else [tags])
            return f"[记忆已保存] ID={mid}"

        async def recall_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            query = params.get("query", "")
            results = self._memory.search(query)
            if not results:
                return "[回忆] 未找到相关记忆"
            lines = [f"- [{r.id}] {r.content[:60]}... (相关度: {r.score:.2f})" for r in results]
            return "[回忆结果]\n" + "\n".join(lines)

        async def summarize_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            period = params.get("period", "today")
            results = self._memory.summarize(period=period)
            return f"[摘要] 共 {len(results)} 条记忆"

        async def list_summaries_tool(**params):
            if self._memory is None:
                return "错误：记忆管理器未初始化"
            limit = params.get("limit", 10)
            return self._memory.list(limit=limit)

        self._tools.register("remember", "保存记忆", remember_tool)
        self._tools.register("recall", "回忆记忆", recall_tool)
        self._tools.register("summarize", "生成摘要", summarize_tool)
        self._tools.register("list_summaries", "列出记忆摘要", list_summaries_tool)

        # ─── 子代理结果查询工具 ───
        async def get_subagent_result_tool(**params):
            if self._subagent_manager is None:
                return "错误：子代理管理器未初始化"
            task_id = params.get("task_id", "")
            task = self._subagent_manager.check_task(task_id)
            if not task:
                return f"任务 {task_id} 不存在"
            if task.status.value != "done":
                return f"任务 {task_id} 状态: {task.status.value}"
            return f"[子代理结果] {task.result.output[:500]}"

        self._tools.register("get_subagent_result", "获取子代理结果", get_subagent_result_tool)

    # ============================================================
    # 系统提示词
    # ============================================================

    def _build_system_prompt(self) -> str:
        """构建系统提示词（包含所有可用工具）"""
        tool_desc_list = []
        for name, info in self._tools._tools.items():
            tool_desc_list.append(f"- {name}: {info['description']}")
        tools_block = "\n".join(tool_desc_list)

        return f"""你是仙人掌 Agent（XianRenZhang Agent），一个自主 AI 助手，支持多平台 LLM。

=== 核心指令 ===
当用户提出任务时，你必须根据需要使用工具来完成任务。工具调用格式：

@@@@
{{
    "tool": "工具名称",
    "params": {{...}},
    "id": "唯一标识"
}}
@@@@

RAW 命令格式（直接执行shell）：
<<<RAW>>>
命令内容
<<<RAW>>>

=== 可用工具 ===
{tools_block}

=== 多平台 LLM 支持 ===
你可以使用以下平台（需要登录）：
- DeepSeek（默认平台）
- 通义千问（tongyi.aliyun.com）
- 豆包（doubao.com）
- 元宝（yuanbao.tencent.com）
- ChatGPT（chat.openai.com）
- Gemini（gemini.google.com）
- Ollama（本地模型 qwen3.5:0.8b）

不同平台的特性：
- DeepSeek: 支持深度思考、文件上传
- ChatGPT: 支持深度思考（Analysis/Extended 模式）、文件上传
- Gemini: 支持深度思考（Think 模式）
- 通义千问、豆包、元宝：基本对话 + 文件上传
- Ollama: 完全本地运行，无需登录

=== 关键规则 ===
1. **禁止直接在 DeepSeek 界面搜索** - 所有搜索必须使用 `browser_search` 工具
2. 一次只能执行一个工具调用
3. 工具执行后，系统会返回结果，你需要根据结果决定下一步
4. 如需多步骤，请分多次发送工具调用
5. 任务完成后，发送 done() 表示结束
6. 如需用户输入，使用 ask("问题")

=== 子代理规则 ===
- browser_research: 启动独立研究子代理（后台运行，不阻塞母代理）
- browser_visit: 启动独立访问子代理
- 子代理启动后返回 task_id，使用 check_task(task_id) 查询进度
- 子代理不能再创建子代理（MAX_DEPTH=1）

=== 记忆系统 ===
- remember(content, tags): 保存重要信息到长期记忆
- recall(query): 搜索相关记忆
- 每运行约 20 轮会自动提醒保存记忆

=== 会话历史追溯 ===
- 每次对话会生成唯一的 DeepSeek URL，可通过该 URL 追溯历史
- 使用 save_conversation() 保存当前对话上下文
- 使用 load_conversation(url=...) 加载历史对话

=== 工作目录 ===
所有文件操作默认相对于: {self._work_dir}
"""

    # ============================================================
    # 公共 API
    # ============================================================

    async def start(self, session: Optional[DeepSeekSession] = None):
        """启动 Commander（复用外部 session）"""
        if session:
            self._session = session
        self._running = True
        # 设置系统提示词到 session，这样每轮对话都会包含
        if self._session:
            self._session.set_system_prompt(self._system_prompt)

    async def run(self, user_instruction: str, file_path: Optional[str] = None,
                  context_hints: str = "") -> str:
        """
        单轮执行（自动循环直到 AI 认为完成）
        返回最终回复内容
        """
        if self._session is None:
            raise RuntimeError("Session 未初始化")

        # 构建第一轮输入（包含系统提示词只在这一轮，后续已存入session）
        current_input = user_instruction
        if file_path:
            current_input += f"\n\n[参考文件: {file_path}]"
        if context_hints:
            current_input += f"\n\n[上下文提示: {context_hints}]"

        final_reply = ""
        max_turns, turn = 1000000000000000, 0

        while turn < max_turns and self._running:
            turn += 1
            self._history_turns += 1

            # 记忆提醒
            if self._memory and self._history_turns % 10 == 0:
                self._emit(EventType.MEMORY_REMINDER, {"turn": self._history_turns})

            # 发送给 AI
            try:
                response = await self._session.send(current_input)
            except Exception as e:
                logger.error(f"AI 调用失败: {e}")
                return f"[错误] AI 调用失败: {e}"

            # 解析指令
            cmds = self._protocol.extract_all(response)
            if not cmds:
                ai_text = response.strip()
                if ai_text:
                    current_input = (
                        "[SYSTEM] NO @@@@ PROTOCOL. Your text was: " + ai_text[:200] + ". All responses MUST use @@@@ JSON format. Format: @@@@\\n{\"tool\":\"xxx\",\"params\":{},\"id\":\"1\"}\\n@@@@. Please retry."
                    )
                    continue
                else:
                    final_reply = response
                    break

            # 执行第一个指令
            cmd = cmds[0]
            self._emit(EventType.COMMAND_DETECTED, {"tool": cmd.command.tool, "id": cmd.id})
            self._emit(EventType.TOOL_START, {"tool": cmd.command.tool})

            result = await self._execute_command(cmd.command)

            self._emit(EventType.TOOL_END, {"tool": cmd.command.tool, "status": result.status})
            self._emit(EventType.COMMAND_SUCCESS if result.status == "success" else EventType.COMMAND_ERROR, result)

            # 构建下一轮输入（只包含工具结果，系统提示词已在session中）
            current_input = f"[系统] 工具 {result.tool} 执行结果:\n{result.output or result.error}\n\n请继续。"

        return final_reply if final_reply else "[完成]"

    async def run_with_loop(self, user_instruction: str, file_path: Optional[str] = None,
                            context_hints: str = "") -> str:
        """带自动循环的 run（和 run 相同，但名称更清晰）"""
        return await self.run(user_instruction, file_path, context_hints)

    async def continue_dialog(self) -> str:
        """
        继续对话（用户发送 '继续' 时调用）
        返回 AI 的回复内容
        """
        if self._correction_pending:
            correction = self._correction_pending
            self._correction_pending = None
            return await self.run(f"[系统纠正]\n{correction}")

        # 发送继续指令
        return await self.run("[系统] 用户要求继续，请接着上一步执行。")

    async def remember(self, content: str, tags: Optional[list] = None):
        """显式保存记忆"""
        if self._memory:
            mid = self._memory.save(content, tags=tags or [])
            return mid
        return None

    async def recall(self, query: str) -> list:
        """显式回忆记忆"""
        if self._memory:
            return self._memory.search(query)
        return []

    def set_memory_manager(self, mm: MemoryManager):
        """设置记忆管理器（外部注入）"""
        self._memory = mm

    def set_subagent_manager(self, sm):
        """设置子代理管理器（外部注入）"""
        self._subagent_manager = sm

    def stop(self):
        """停止运行循环"""
        self._running = False

    def inject_correction(self, correction_text: str):
        """注入系统纠正（AI 幻觉时人工干预）"""
        self._correction_pending = correction_text

    # ============================================================
    # 内部方法
    # ============================================================

    async def _execute_command(self, cmd: SimpleNamespace) -> ExecutionResult:
        """执行单个命令"""
        tool_name = getattr(cmd, "tool", None) or getattr(cmd, "action", None)
        params = getattr(cmd, "params", {}) or {}
        cmd_id = getattr(cmd, "id", str(uuid.uuid4()))

        if not tool_name:
            return ExecutionResult(id=cmd_id, status="error", error="指令缺少 tool 字段")

        # 特殊指令处理
        if tool_name == "done":
            return ExecutionResult(id=cmd_id, status="success", output="任务完成", tool="done")

        if tool_name == "ask":
            question = params.get("question", "请输入：")
            return ExecutionResult(id=cmd_id, status="success", output=f"ASK:{question}", tool="ask")

        # RAW 命令处理
        if tool_name == "raw_shell":
            return await self._tools.execute("shell_exec", {"command": params.get("command", "")})

        # 标准工具调用
        return await self._tools.execute(tool_name, params)

    def _emit(self, event_type: EventType, data=None):
        """触发事件回调"""
        if self._on_event:
            try:
                self._on_event(AgentEvent(event_type=event_type, data=data))
            except Exception as e:
                logger.warning(f"事件回调错误: {e}")

    # ============================================================
    # 错误纠正（commander fix）
    # ============================================================

    async def fix(self, error_hint: str):
        """
        当检测到 AI 幻觉/错误时，注入系统级纠正
        使用方式：在 terminal.py 检测到异常后调用 commander.fix("纠正内容")
        """
        self.inject_correction(error_hint)
        return await self.continue_dialog()


# ============================================================
# 快捷函数（兼容旧代码）
# ============================================================

async def run_agent(browser_manager, session: DeepSeekSession, instruction: str, work_dir: str = "") -> str:
    """快捷函数：创建 Commander 并执行"""
    commander = Commander(
        browser_manager=browser_manager,
        session=session,
        work_dir=work_dir or os.getcwd(),
    )
    await commander.start(session=session)
    return await commander.run(instruction)
